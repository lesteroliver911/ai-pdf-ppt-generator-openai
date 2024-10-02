import io
import json
import logging
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from langchain_openai import ChatOpenAI
from langchain_core.prompts import PromptTemplate
from utils.image_handling import search_pexels_images, add_image_to_slide, add_credits_slide

logger = logging.getLogger(__name__)

LLM_MODEL_OPENAI: str = os.getenv('LLM_MODEL_OPENAI', 'gpt-4o')

def clean_json_string(content: str) -> str:
    # Remove any non-printable characters (including control characters)
    content = ''.join(char for char in content if char.isprintable())
    # Remove any leading/trailing whitespace
    content = content.strip()
    # Remove ```json at the start and ``` at the end if present
    content = re.sub(r'^```json\s*', '', content)
    content = re.sub(r'\s*```$', '', content)
    # Replace any sequence of newlines and spaces with a single space
    content = re.sub(r'\s+', ' ', content)
    return content

def parse_json_safely(content: str) -> dict:
    cleaned_content = clean_json_string(content)
    try:
        return json.loads(cleaned_content)
    except json.JSONDecodeError as e:
        logger.error(f"JSON Decode Error: {str(e)}")
        logger.error(f"Cleaned Content: {cleaned_content}")
        raise ValueError("Failed to parse AI response. Please try again.")

def generate_presentation_content(topic: str, information: str, user_instructions: str, num_slides: int) -> dict:
    prompt = PromptTemplate(
        input_variables=["topic", "information", "instructions", "num_slides"],
        template="""
    You are an expert Power Point presentation creator. Create a compelling presentation about {topic} with exactly {num_slides} slides, Mandatory.
    Use the following information as a basis, but feel free to expand on it with your knowledge:
    ---
    {information}
    ---
    User Instructions:
    {instructions}

    Structure the presentation according to the user instructions and the specified number of slides. Always include:
    1. An engaging introduction slide
    2. A conclusion slide summarizing the main points
    3. A "Next Steps" or "Call to Action" slide that includes the contact details if any, !imporant

    Adjust the content to fit the specified number of slides. If more slides are requested, expand on the topic with more details.
    If fewer slides are requested, focus on the most important points and summarize the content accordingly.

    Each slide should have:
    - A concise, attention-grabbing title
    - Points of key information (written in full sentences), elaborate if required
    - A brief speaker note providing additional context or talking points.

    Return the structured information as a JSON as follows:
    {{
        "slides": [
            {{
                "title": "Slide Title",
                "content": "Bullet point 1\\n Bullet point 2\\n Bullet point 3",
                "speaker_note": "Additional context or talking points for the presenter"
            }},
            // ... more slides ...
        ]
    }}
    """
    )
    
    llm = ChatOpenAI(model=LLM_MODEL_OPENAI, temperature=0.5)
    response = llm.invoke(prompt.format(topic=topic, information=information, instructions=user_instructions, num_slides=num_slides))
    
    return parse_json_safely(response.content)

def customize_presentation_style(slides: dict, style: str, audience: str) -> dict:
    prompt = PromptTemplate(
        input_variables=["slides", "style", "audience"],
        template="""
    Adapt the following presentation slides to a {style} style for a {audience} audience. Mandatory.
    Maintain the overall structure and key points, but adjust the language, tone, and complexity accordingly.

    Original slides:
    {slides}

    Return the adapted slides in the same JSON format. Do NOT HALLUCINATE. Make sure the JSON format is passed. Nothing else, Take a deep breath before starting any task.
    """
    )
    
    llm = ChatOpenAI(model=LLM_MODEL_OPENAI, temperature=0.5)
    response = llm.invoke(prompt.format(slides=json.dumps(slides), style=style, audience=audience))
    
    return parse_json_safely(response.content)

def generate_executive_summary(slides: dict) -> dict:
    prompt = PromptTemplate(
        input_variables=["slides"],
        template="""
    Create a concise executive summary for the following presentation:
    {slides}

    The summary should:
    1. Capture the main idea of the presentation in one sentence
    2. List Key takeaways
    3. Conclude with a brief statement of the presentation's significance or call to action

    Return the summary as a JSON object with the following structure:
    {{
        "title": "Executive Summary",
        "content": " Main idea\n Key takeaway 1\n Key takeaway 2\n Key takeaway 3\n Conclusion",
        "speaker_note": "Additional context or emphasis points"
    }}
    """
    )
    
    llm = ChatOpenAI(model=LLM_MODEL_OPENAI, temperature=0.5)
    response = llm.invoke(prompt.format(slides=json.dumps(slides)))
    
    return parse_json_safely(response.content)

# The rest of the functions (apply_template_styling, replace_fonts, create_presentation) remain unchanged

def apply_template_styling(slide, template_choice):
    if template_choice == 'dark_modern':
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(255, 255, 255)  # White text
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
    elif template_choice == 'bright_modern':
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

def replace_fonts(presentation):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
    logger.info("Replaced fonts in all slides")

def create_presentation(slides: dict, template_choice: str, presentation_title: str, presenter_name: str, insert_image: bool, num_slides: int) -> io.BytesIO:
    prs = Presentation()

    # Add title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    title.text = presentation_title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Presented by {presenter_name}"
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    apply_template_styling(title_slide, template_choice)
    logger.info("Added title slide")

    # Calculate the number of content slides to add
    content_slides_to_add = num_slides - 1  # Subtract 1 for the title slide 

    # Add content slides
    for i, slide in enumerate(slides['slides'], 1):
        if i > content_slides_to_add:
            break

        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = content_slide.shapes.title
        title.text = slide['title']
        title.text_frame.paragraphs[0].font.size = Pt(36)
        content = content_slide.placeholders[1]
        content.text = slide['content']
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
        
        if 'speaker_note' in slide:
            content_slide.notes_slide.notes_text_frame.text = slide['speaker_note']
        
        apply_template_styling(content_slide, template_choice)
        logger.info(f"Added content slide {i}: {slide['title']}")

        if insert_image:
            add_image_to_slide(content_slide, slide['title'])

    if insert_image:
        # Only add credits slide if there's room and images were inserted
        if len(prs.slides) < num_slides:
            add_credits_slide(prs, template_choice)

    replace_fonts(prs)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    logger.info(f"Presentation has been created with {len(prs.slides)} slides")
    return pptx_io
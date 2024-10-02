import io
import requests
import logging
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from config import PEXELS_API_KEY

logger = logging.getLogger(__name__)

def search_pexels_images(keyword):
    query = keyword.lower().replace(' ', '%20')
    PEXELS_API_URL = f'https://api.pexels.com/v1/search?query={query}&per_page=1'
    headers = {
        'Authorization': PEXELS_API_KEY
    }
    try:
        response = requests.get(PEXELS_API_URL, headers=headers)
        response.raise_for_status()
        data = response.json()
        if 'photos' in data and len(data['photos']) > 0:
            return data['photos'][0]['src']['medium']
    except requests.RequestException as e:
        logger.error(f"Error fetching image from Pexels: {str(e)}")
    return None

def add_image_to_slide(slide, keyword):
    image_url = search_pexels_images(keyword)
    if image_url:
        try:
            image_data = requests.get(image_url).content
            image_stream = io.BytesIO(image_data)
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            image_width = Inches(4)
            image_height = Inches(3)
            left = slide_width - image_width - Inches(0.5)
            top = slide_height - image_height - Inches(0.5)
            slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)
            logger.info(f"Added image for keyword: {keyword}")
        except Exception as e:
            logger.error(f"Error adding image to slide: {str(e)}")
    else:
        logger.warning(f"No image found for keyword: {keyword}")

def add_credits_slide(prs, template_choice):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Credits"
    title_frame.paragraphs[0].font.size = Pt(44)
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(1)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame
    paragraph = content_frame.add_paragraph()
    paragraph.text = "Images provided by Pexels: https://www.pexels.com"
    
    if template_choice == 'dark_modern':
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(255, 255, 255)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    elif template_choice == 'bright_modern':
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    logger.info("Added credits slide")
